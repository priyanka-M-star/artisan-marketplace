from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# Set slide size to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=None):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Set content
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.clear()

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(10)
        if item.startswith("  "):
            p.level = 1
            p.text = item.strip()

    return slide

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide = add_title_slide(prs, "Artisan Marketplace", "A Platform for Handcrafted Goods\nProject Presentation")

# ============================================================================
# SLIDE 2: Introduction
# ============================================================================
slide = add_content_slide(prs, "Introduction", [
    "Artisan Marketplace is a web-based e-commerce platform",
    "Connects skilled artisans directly with customers",
    "Eliminates middlemen to ensure fair pricing",
    "Supports multiple product categories:",
    "  - Pottery, Jewelry, Paintings, Crafts, Textiles",
    "Built with modern MERN stack architecture"
])

# ============================================================================
# SLIDE 3: Literature Survey
# ============================================================================
slide = add_content_slide(prs, "Literature Survey", [
    "Existing platforms: Etsy, Amazon Handmade, IndiaMART",
    "Limitations identified in existing systems:",
    "  - High commission fees (15-25%)",
    "  - Limited customization for artisans",
    "  - Complex vendor onboarding",
    "  - Lack of integrated payment solutions",
    "Our solution addresses these gaps with:",
    "  - Lower commissions via Stripe Connect",
    "  - Role-based access control",
    "  - Simplified vendor dashboard"
])

# ============================================================================
# SLIDE 4: Week 1 - Project Setup & Database Design
# ============================================================================
slide = add_content_slide(prs, "Week 1: Project Setup & Database Design", [
    "Initialized Node.js + Express backend",
    "Configured MongoDB database connection",
    "Created database schemas:",
    "  - User (buyers with OTP verification)",
    "  - Vendor (artisans with Stripe accounts)",
    "  - Product (with image uploads)",
    "  - Order (tracking & status)",
    "Set up project structure with MVC pattern",
    "Configured environment variables & security (Helmet, CORS)"
])

# ============================================================================
# SLIDE 5: Week 2 - Authentication & Authorization
# ============================================================================
slide = add_content_slide(prs, "Week 2: Authentication & Authorization", [
    "Implemented JWT-based authentication",
    "Built registration with role selection (Buyer/Vendor)",
    "Password strength validation (uppercase, lowercase, number, symbol)",
    "OTP-based password reset functionality",
    "Created authentication middleware for route protection",
    "Developed login page with responsive UI",
    "Session management using localStorage"
])

# ============================================================================
# SLIDE 6: Week 3 - Product & Vendor Features
# ============================================================================
slide = add_content_slide(prs, "Week 3: Product & Vendor Features", [
    "Product CRUD APIs (Create, Read, Update, Delete)",
    "Image upload integration with Cloudinary",
    "Product search and pagination",
    "Vendor dashboard for managing products",
    "Shipping profile configuration",
    "Admin routes for platform management",
    "Role guards for access control"
])

# ============================================================================
# SLIDE 7: Week 4 - Orders & Payment Integration
# ============================================================================
slide = add_content_slide(prs, "Week 4: Orders & Payment Integration", [
    "Order management APIs (create, view, cancel)",
    "Stripe Connect integration for vendor payouts",
    "Commission splitting between platform & vendors",
    "Webhook handling for payment events",
    "Shipping routes and configuration",
    "Payment routes and controller",
    "Frontend checkout page implementation"
])

# ============================================================================
# SLIDE 8: System Architecture
# ============================================================================
slide = add_content_slide(prs, "System Architecture", [
    "Frontend: HTML5, CSS3, Vanilla JavaScript",
    "Backend: Node.js + Express.js",
    "Database: MongoDB with Mongoose ODM",
    "Cloud Services:",
    "  - Cloudinary (image storage)",
    "  - Stripe (payment processing)",
    "Security: Helmet, bcryptjs, JWT",
    "Architecture: REST API with MVC pattern"
])

# ============================================================================
# SLIDE 9: Conclusion
# ============================================================================
slide = add_content_slide(prs, "Conclusion", [
    "Successfully built a functional marketplace platform",
    "Key achievements:",
    "  - Complete user authentication system",
    "  - Vendor onboarding with Stripe Connect",
    "  - Product management with image uploads",
    "  - Order tracking and cancellation",
    "  - Secure payment processing",
    "Future enhancements:",
    "  - Email notifications",
    "  - Product reviews & ratings",
    "  - Mobile application"
])

# ============================================================================
# SLIDE 10: References
# ============================================================================
slide = add_content_slide(prs, "References", [
    "Express.js Documentation - https://expressjs.com/",
    "MongoDB & Mongoose Docs - https://mongodb.com/",
    "Stripe Connect API - https://stripe.com/docs/connect",
    "Cloudinary Documentation - https://cloudinary.com/docs",
    "JWT Authentication Best Practices - https://jwt.io/",
    "Etsy Seller Handbook - Market research reference",
    "MDN Web Docs - https://developer.mozilla.org/"
])

# Save the presentation
output_file = r"C:\Users\priya\OneDrive\Desktop\artisan-marketplace\Artisan_Marketplace_Presentation.pptx"
prs.save(output_file)

print(f"Presentation saved to: {output_file}")
print("10 slides created successfully!")
